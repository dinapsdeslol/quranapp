from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height

BLUE_DARK = RGBColor(0x0D, 0x47, 0xA1)
BLUE_MID = RGBColor(0x15, 0x65, 0xC0)
BLUE_LIGHT = RGBColor(0x19, 0x76, 0xD2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
WHITE70 = RGBColor(0xB0, 0xB0, 0xB0)
BLACK = RGBColor(0x00, 0x00, 0x00)
GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
GREEN = RGBColor(0x4C, 0xAF, 0x50)
RED = RGBColor(0xF4, 0x43, 0x36)
YELLOW = RGBColor(0xFF, 0xC1, 0x07)
CODE_BG = RGBColor(0x1E, 0x1E, 0x2E)
CODE_TXT = RGBColor(0xCD, 0xD6, 0xF4)

def add_bg(slide, color=BLUE_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_rounded_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_tb(slide, left, top, width, height, text, size=18, color=BLACK, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = 'Calibri'
    p.alignment = align
    return txBox

def add_para(tf, text, size=16, color=BLACK, bold=False, align=PP_ALIGN.LEFT, sp=Pt(6), sa=Pt(4)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = 'Calibri'
    p.alignment = align
    p.space_before = sp
    p.space_after = sa
    return p

def add_bul(tf, text, size=14, color=BLACK, bold=False):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = 'Calibri'
    p.level = 0
    p.space_before = Pt(2)
    p.space_after = Pt(2)
    return p

def new_slide():
    return prs.slides.add_slide(prs.slide_layouts[6])

def code_block(slide, left, top, width, height, code_text, size=10):
    shape = add_rounded_rect(slide, left, top, width, height, CODE_BG)
    txBox = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.1), width - Inches(0.3), height - Inches(0.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    lines = code_text.split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = CODE_TXT
        p.font.name = 'Consolas'
        p.space_before = Pt(0)
        p.space_after = Pt(0)
    return txBox

# ============================================================
# HELPERS FOR STANDARD SLIDES
# ============================================================
def title_slide(text):
    slide = new_slide()
    add_bg(slide, WHITE)
    add_shape(slide, 0, 0, W, Inches(1.1), BLUE_DARK)
    add_tb(slide, Inches(0.8), Inches(0.2), Inches(11), Inches(0.8), text, 36, WHITE, True)
    return slide

def add_subtitle_box(slide, y, title, items, title_size=20, item_size=14):
    add_tb(slide, Inches(0.5), Inches(y), Inches(12), Inches(0.4), title, title_size, BLUE_DARK, True)
    for i, item in enumerate(items):
        add_tb(slide, Inches(0.8), Inches(y + 0.4 + i * 0.38), Inches(12), Inches(0.36), item, item_size, GRAY)

# ============================================================
# SLIDE 1: TITLE
# ============================================================
slide = new_slide()
add_bg(slide, BLUE_DARK)
add_shape(slide, 0, Inches(3.2), W, Inches(0.06), BLUE_LIGHT)
add_tb(slide, Inches(1), Inches(1.2), Inches(11), Inches(1.5), 'Quran Player', 60, WHITE, True, PP_ALIGN.CENTER)
add_tb(slide, Inches(1), Inches(2.6), Inches(11), Inches(0.8), 'Application Mobile de Lecture Audio du Coran', 28, WHITE70, align=PP_ALIGN.CENTER)
add_tb(slide, Inches(1), Inches(3.6), Inches(11), Inches(0.6), 'ING3 - Sécurité du Développement Mobile', 22, WHITE70, align=PP_ALIGN.CENTER)
add_tb(slide, Inches(1), Inches(4.5), Inches(11), Inches(0.5), 'Flutter  |  Firebase  |  Authentification Biométrique  |  Audio Streaming', 18, WHITE70, align=PP_ALIGN.CENTER)
add_tb(slide, Inches(1), Inches(5.3), Inches(11), Inches(0.5), 'Développé par : [Votre Nom]  |  2025/2026', 16, WHITE70, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 2: PLAN
# ============================================================
slide = title_slide('Plan de la Présentation')
items = [('1. Introduction', 'Contexte, objectifs, technologies choisies'),
    ('2. Outils & Langages', 'Flutter, Dart, Firebase, Android Studio, Git, libraries utilisées'),
    ('3. Architecture Générale', 'Structure MVC du projet, arborescence lib/, organisation des dossiers'),
    ('4. main.dart - Point d\'Entrée', 'Initialisation Firebase, AppFlow, machine à états, routing conditionnel'),
    ('5. bio_screen.dart - Biométrie', 'local_auth, isFirstLaunch, authentification fingerprint, skip'),
    ('6. auth_service.dart - Firebase Auth', 'Login, Register, ResetPassword, getUserData, fallback Firestore'),
    ('7. login_screen.dart & register_screen.dart', 'Formulaires, validation, âge ≥13, Firebase Auth'),
    ('8. home_screen.dart - Navigation', 'IndexedStack, AudioService partagé, logout avec cleanup'),
    ('9. player_screen.dart - Lecture Audio', 'API Quran, audioplayers, contrôles, favoris, slider'),
    ('10. stats_screen.dart - Statistiques', 'fl_chart, SharedPreferences, objectif mensuel, top tracks'),
    ('11. favorites_screen.dart - Firestore', 'StreamBuilder, CRUD, biométrie pour delete, SnackBar erreurs'),
    ('12. audio_service.dart - Coeur Audio', 'playTrack, togglePlayPause, enregistrement stats, gestion erreurs'),
    ('13. api_service.dart & Modèles', 'staticquran API, fallback, AudioTrack, SurahCategory'),
    ('14. Sécurité & Difficultés', 'Biométrie, permissions, Firestore rules, problèmes rencontrés'),
    ('15. Conclusion & Demo', 'Résumé, améliorations possibles, démonstration')]
y = 1.3
for title, desc in items:
    add_rounded_rect(slide, Inches(0.5), Inches(y), Inches(12.3), Inches(0.38), LIGHT_GRAY)
    add_tb(slide, Inches(0.8), Inches(y + 0.02), Inches(4), Inches(0.34), title, 15, BLUE_DARK, True)
    add_tb(slide, Inches(5), Inches(y + 0.02), Inches(7.5), Inches(0.34), desc, 13, GRAY)
    y += 0.41

# ============================================================
# SLIDE 3: INTRODUCTION
# ============================================================
slide = title_slide('1. Introduction')
tb = add_tb(slide, Inches(0.5), Inches(1.3), Inches(12), Inches(5.5), '', 14, BLACK)
tf = tb.text_frame; tf.word_wrap = True
tf.paragraphs[0].text = 'Contexte du Projet'; tf.paragraphs[0].font.size = Pt(24); tf.paragraphs[0].font.bold = True; tf.paragraphs[0].font.color.rgb = BLUE_DARK
add_bul(tf, 'Application mobile Android développée dans le cadre du cours ING3 - Sécurité du Développement Mobile', 15)
add_bul(tf, 'Objectif : créer une application sécurisée de lecture audio du Coran avec authentification biométrique', 15)
add_bul(tf, 'Technologies : Flutter (cross-platform), Firebase (backend), local_auth (biométrie)', 15)
add_para(tf, '', 6)
add_para(tf, 'Fonctionnalités Clés', 22, BLUE_DARK, True)
add_bul(tf, 'Authentification biométrique (fingerprint) au premier lancement - protection de l\'accès à l\'app', 15)
add_bul(tf, 'Firebase Authentication pour la gestion des comptes (email/mot de passe)', 15)
add_bul(tf, 'Cloud Firestore pour la synchronisation des favoris en temps réel', 15)
add_bul(tf, 'Lecture audio des 114 sourates du Coran via API REST + streaming MP3', 15)
add_bul(tf, 'Statistiques d\'écoute (temps total, histogramme 30 jours, top tracks, objectif mensuel)', 15)
add_bul(tf, 'Architecture MVC (Model-View-Controller) propre et maintenable', 15)
add_para(tf, '', 6)
add_para(tf, 'Choix Technologiques', 22, BLUE_DARK, True)
add_bul(tf, 'Flutter : développement cross-platform, hot reload, Material 3, riche écosystème de packages', 15)
add_bul(tf, 'Firebase : backend sans serveur, auth sécurisée, base de données temps réel', 15)
add_bul(tf, 'local_auth : abstraction de l\'authentification biométrique Android/iOS', 15)
add_bul(tf, 'audioplayers : lecture audio multi-plateforme avec contrôle du state', 15)
add_bul(tf, 'fl_chart : graphiques statistiques personnalisables pour Flutter', 15)

# ============================================================
# SLIDE 4: OUTILS & LANGAGES
# ============================================================
slide = title_slide('2. Outils & Langages')
# Left column - Tools
add_tb(slide, Inches(0.5), Inches(1.3), Inches(6), Inches(0.4), 'Outils de Développement', 20, BLUE_DARK, True)
tools_left = ['Flutter 3.41.3 - Framework UI cross-platform', 'Dart 3.11.1 - Langage de programmation', 'Android Studio + AVD Manager - IDE & Émulateur', 'Firebase Console - Configuration backend', 'Git / GitHub - Contrôle de version', 'JDK 17 (Eclipse Adoptium Temurin)', 'Gradle 8.11.1 - Build system Android', 'Python 3.14 - Génération de la présentation']
for i, t in enumerate(tools_left):
    add_tb(slide, Inches(0.5), Inches(1.8 + i * 0.35), Inches(6), Inches(0.33), f'  \u2022 {t}', 13, GRAY)

# Right column - Libraries
add_tb(slide, Inches(6.8), Inches(1.3), Inches(6), Inches(0.4), 'Packages Flutter (pubspec.yaml)', 20, BLUE_DARK, True)
libs = ['firebase_core ^3.6.0 - Initialisation Firebase', 'firebase_auth ^5.3.1 - Authentification', 'cloud_firestore ^5.4.4 - Base de données', 'local_auth ^2.3.0 - Biométrie', 'audioplayers ^6.1.0 - Lecture audio', 'http ^1.2.2 - Requêtes HTTP', 'fl_chart ^0.69.0 - Graphiques', 'shared_preferences ^2.3.2 - Stockage local', 'intl ^0.19.0 - Formatage dates']
for i, l in enumerate(libs):
    add_tb(slide, Inches(6.8), Inches(1.8 + i * 0.35), Inches(6), Inches(0.33), f'  \u2022 {l}', 13, GRAY)

# Langages
add_tb(slide, Inches(0.5), Inches(5.2), Inches(12), Inches(0.4), 'Langages de Programmation', 20, BLUE_DARK, True)
langs = [('Dart', '~95% du code', 'Toute la logique : UI, services, modèles, états'),
    ('Kotlin', '~3%', 'Configuration Gradle, build scripts Android'),
    ('XML', '~1.5%', 'AndroidManifest, ressources thème, google-services.json'),
    ('Python', '~0.5%', 'Génération de cette présentation (python-pptx)')]
for i, (n, p, d) in enumerate(langs):
    x = Inches(0.5 + i * 3.1)
    add_rounded_rect(slide, x, Inches(5.7), Inches(2.9), Inches(1.2), LIGHT_GRAY)
    add_tb(slide, x + Inches(0.15), Inches(5.75), Inches(2.6), Inches(0.35), n, 18, BLUE_DARK, True)
    add_tb(slide, x + Inches(0.15), Inches(6.1), Inches(2.6), Inches(0.3), p, 16, GREEN, True)
    add_tb(slide, x + Inches(0.15), Inches(6.4), Inches(2.6), Inches(0.45), d, 12, GRAY)

# ============================================================
# SLIDE 5: ARCHITECTURE GENERALE
# ============================================================
slide = title_slide('3. Architecture Générale du Projet')
add_tb(slide, Inches(0.5), Inches(1.2), Inches(12), Inches(0.4), 'Arborescence lib/', 20, BLUE_DARK, True)
tree_text = 'lib/\n' + '  \u2514\u2500\u2500 main.dart                 # Point d\'entr\u00e9e + AppFlow (machine \u00e0 \u00e9tats)\n' + '  \u2514\u2500\u2500 models/\n' + '       \u2514\u2500\u2500 track_model.dart       # AudioTrack, SurahCategory\n' + '  \u2514\u2500\u2500 services/\n' + '       \u2514\u2500\u2500 auth_service.dart     # Firebase Auth + Firestore user data\n' + '       \u2514\u2500\u2500 bio_service.dart      # local_auth + SharedPreferences first_launch\n' + '       \u2514\u2500\u2500 audio_service.dart    # audioplayers + stats recording\n' + '       \u2514\u2500\u2500 api_service.dart      # HTTP client pour staticquran API\n' + '       \u2514\u2500\u2500 favorite_service.dart # Firestore CRUD pour les favoris\n' + '  \u2514\u2500\u2500 screens/\n' + '       \u2514\u2500\u2500 bio_screen.dart       # \u00c9cran d\'auth biom\u00e9trique\n' + '       \u2514\u2500\u2500 login_screen.dart     # \u00c9cran de connexion\n' + '       \u2514\u2500\u2500 register_screen.dart  # \u00c9cran d\'inscription\n' + '       \u2514\u2500\u2500 forgot_screen.dart    # Mot de passe oubli\u00e9\n' + '       \u2514\u2500\u2500 home_screen.dart      # Navigation principale (3 tabs)\n' + '       \u2514\u2500\u2500 player_screen.dart    # Lecteur audio + liste sourates\n' + '       \u2514\u2500\u2500 stats_screen.dart     # Statistiques d\'\u00e9coute\n' + '       \u2514\u2500\u2500 favorites_screen.dart # Favoris synchronis\u00e9s Firestore'
add_tb(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(5.5), tree_text, 12, GRAY)

# Architecture schema
add_tb(slide, Inches(6.5), Inches(1.2), Inches(6), Inches(0.4), 'Architecture MVC', 20, BLUE_DARK, True)
# Models
add_rounded_rect(slide, Inches(6.8), Inches(1.7), Inches(5.3), Inches(0.55), BLUE_MID)
add_tb(slide, Inches(6.8), Inches(1.72), Inches(5.3), Inches(0.5), 'MOD\u00c8LES (Models)', 16, WHITE, True, PP_ALIGN.CENTER)
add_tb(slide, Inches(6.8), Inches(2.28), Inches(5.3), Inches(0.3), 'track_model.dart (AudioTrack, SurahCategory)', 11, GRAY, align=PP_ALIGN.CENTER)
# Arrow
add_tb(slide, Inches(8.5), Inches(2.6), Inches(2), Inches(0.3), '\u25bc', 16, BLUE_MID, align=PP_ALIGN.CENTER)
# Services
add_rounded_rect(slide, Inches(6.8), Inches(2.9), Inches(5.3), Inches(2.5), LIGHT_GRAY)
add_tb(slide, Inches(6.8), Inches(2.95), Inches(5.3), Inches(0.4), 'SERVICES (Contrôleurs)', 16, BLUE_DARK, True, PP_ALIGN.CENTER)
svcs = ['auth_service.dart : Firebase Auth + Firestore', 'bio_service.dart : local_auth + SharedPrefs', 'audio_service.dart : audioplayers + stats', 'api_service.dart : HTTP (staticquran)', 'favorite_service.dart : Firestore CRUD']
for i, s in enumerate(svcs):
    add_tb(slide, Inches(7), Inches(3.4 + i * 0.38), Inches(5), Inches(0.36), f'  \u2022 {s}', 11, GRAY)
# Arrow
add_tb(slide, Inches(8.5), Inches(5.5), Inches(2), Inches(0.3), '\u25bc', 16, BLUE_MID, align=PP_ALIGN.CENTER)
# Screens
add_rounded_rect(slide, Inches(6.8), Inches(5.8), Inches(5.3), Inches(1.4), LIGHT_GRAY)
add_tb(slide, Inches(6.8), Inches(5.85), Inches(5.3), Inches(0.4), 'VUES (Screens / UI)', 16, BLUE_DARK, True, PP_ALIGN.CENTER)
add_tb(slide, Inches(7), Inches(6.25), Inches(5), Inches(0.36), '  8 écrans (bio, login, register, forgot, home, player, stats, favorites)', 11, GRAY)
add_tb(slide, Inches(7), Inches(6.55), Inches(5), Inches(0.36), '  main.dart : point d\'entrée + AppFlow (orchestrateur)', 11, GRAY)
# External
add_rounded_rect(slide, Inches(6.8), Inches(1.7), Inches(2.4), Inches(0.55), RGBColor(0x00, 0x96, 0x88))
add_tb(slide, Inches(6.8), Inches(1.72), Inches(2.4), Inches(0.5), 'EXTERNE', 16, WHITE, True, PP_ALIGN.CENTER)
add_tb(slide, Inches(9.5), Inches(1.75), Inches(2.6), Inches(0.5), 'Firebase Auth/Firestore\nstaticquran API\nMP3 Server', 11, GRAY)

# ============================================================
# SLIDE 6: main.dart
# ============================================================
slide = title_slide('4. main.dart - Point d\'Entrée & Machine à États')
tb = add_tb(slide, Inches(0.5), Inches(1.3), Inches(6), Inches(5.5), '', 14, BLACK)
tf = tb.text_frame; tf.word_wrap = True
tf.paragraphs[0].text = 'Fonction main()'; tf.paragraphs[0].font.size = Pt(20); tf.paragraphs[0].font.bold = True; tf.paragraphs[0].font.color.rgb = BLUE_DARK
add_bul(tf, 'WidgetsFlutterBinding.ensureInitialized() - Initialise le binding Flutter', 14)
add_bul(tf, 'Firebase.initializeApp() - Initialise Firebase (uniquement sur mobile, pas web)', 14)
add_bul(tf, 'runApp(QuranPlayerApp()) - Lance l\'application', 14)
add_para(tf, '', 6)
add_para(tf, 'QuranPlayerApp (StatelessWidget)', 20, BLUE_DARK, True)
add_bul(tf, 'MaterialApp avec thème Material 3, seed color #0D47A1 (bleu)', 14)
add_bul(tf, 'debugShowCheckedModeBanner: false - Cache le bandeau DEBUG', 14)
add_bul(tf, 'home: AppFlow() - Le widget racine qui gère le flux', 14)
add_para(tf, '', 6)
add_para(tf, 'AppFlow (StatefulWidget) - Machine à États', 20, BLUE_DARK, True)
add_bul(tf, '4 états possibles : _showBio / _showLogin / _loggedIn + _userData / spinner', 14)
add_bul(tf, '_checkLogin() :', 14, BLUE_DARK, True)
add_bul(tf, '  1. Vérifie isFirstLaunch() via SharedPreferences (clé bio_done)', 13)
add_bul(tf, '  2. Si premier lancement -> _showBio = true (BiometricScreen)', 13)
add_bul(tf, '  3. Sinon, vérifie currentUser (Firebase Auth)', 13)
add_bul(tf, '  4. Si user connecté -> getUserData() depuis Firestore', 13)
add_bul(tf, '  5. Si data OK -> HomeScreen, sinon LoginScreen', 13)
add_para(tf, '', 6)
add_para(tf, 'Callbacks', 20, BLUE_DARK, True)
add_bul(tf, '_onBioSuccess() : biometric OK -> passe au login', 14)
add_bul(tf, '_onLogin(data) : login OK -> HomeScreen avec userData', 14)
add_bul(tf, 'Build() : Affiche l\'écran correspondant à l\'état courant', 14)
# Code block
code = '''class _AppFlowState extends State<AppFlow> {
  Future<void> _checkLogin() async {
    final isFirst = await _bio!.isFirstLaunch();
    if (isFirst) { setState(() => _showBio = true); return; }
    final user = _auth!.currentUser;
    if (user != null) {
      final data = await _auth!.getUserData(user.uid);
      if (data != null) {
        setState(() { _userData = data; _loggedIn = true; });
        return;
      }
    }
    setState(() => _showLogin = true);
  }
  @override
  Widget build(BuildContext context) {
    if (_showBio) return BiometricScreen(onSuccess: _onBioSuccess);
    if (_showLogin) return LoginScreen(onLogin: _onLogin);
    if (_loggedIn && _userData != null) return HomeScreen(userData: _userData!);
    return const Scaffold(body: Center(child: CircularProgressIndicator()));
  }
}'''
code_block(slide, Inches(6.8), Inches(1.3), Inches(6), Inches(5.8), code, 11)

# ============================================================
# SLIDE 7: bio_screen.dart
# ============================================================
slide = title_slide('5. bio_screen.dart - Authentification Biométrique')
tb = add_tb(slide, Inches(0.5), Inches(1.2), Inches(6.2), Inches(5.8), '', 14, BLACK)
tf = tb.text_frame; tf.word_wrap = True
tf.paragraphs[0].text = 'Objectif'; tf.paragraphs[0].font.size = Pt(20); tf.paragraphs[0].font.bold = True; tf.paragraphs[0].font.color.rgb = BLUE_DARK
add_bul(tf, 'Protéger l\'accès à l\'application par empreinte digitale au premier lancement', 14)
add_bul(tf, 'Bouton Skip disponible si le capteur biométrique n\'est pas disponible', 14)
add_para(tf, '', 6)
add_para(tf, 'Fonctionnement Détaillé', 22, BLUE_DARK, True)
add_bul(tf, '_init() : Appelée dans initState(), vérifie la disponibilité du capteur', 14)
add_bul(tf, '  -> BioService.isBiometricAvailable() utilise local_auth', 13)
add_bul(tf, '  -> Si disponible : "Place your finger on the sensor" + _tryAuth()', 13)
add_bul(tf, '  -> Sinon : "No fingerprint registered" + boutons Settings/Skip', 13)
add_bul(tf, '_tryAuth() : Appelle BioService.authenticate()', 14)
add_bul(tf, '  -> local_auth avec AuthenticationOptions(biometricOnly: true)', 13)
add_bul(tf, '  -> Succès : SystemSound.play(SystemSoundType.click) + setFirstLaunchComplete() + onSuccess()', 13)
add_bul(tf, '  -> Échec : message "Authentication failed" + bouton Retry', 13)
add_bul(tf, 'Bouton Skip : setFirstLaunchComplete() + onSuccess() sans auth', 14)
add_para(tf, '', 6)
add_para(tf, 'UI / UX', 20, BLUE_DARK, True)
add_bul(tf, 'Fond dégradé bleu (#0D47A1 -> #1976D2)', 14)
add_bul(tf, 'Icône fingerprint (Icons.fingerprint) en grand (size: 120)', 14)
add_bul(tf, 'Message dynamique selon l\'état du capteur', 14)
add_bul(tf, 'Boutons : Settings (va aux paramètres), Retry (réessaie), Skip (passe)', 14)

code = '''class _BiometricScreenState extends State<BiometricScreen> {
  Future<void> _init() async {
    final available = await _bio.isBiometricAvailable();
    if (!available) {
      setState(() { _msg = 'No fingerprint registered.'; _hasBio = false; });
      return;
    }
    setState(() { _hasBio = true; _msg = 'Place your finger on the sensor'; });
    await _tryAuth();
  }
  Future<void> _tryAuth() async {
    final ok = await _bio.authenticate();
    if (ok) {
      await SystemSound.play(SystemSoundType.click);
      await _bio.setFirstLaunchComplete();
      widget.onSuccess();
    }
  }
  // Skip button:
  onPressed: () {
    _bio.setFirstLaunchComplete();
    widget.onSuccess();
  }
}'''
code_block(slide, Inches(6.8), Inches(1.2), Inches(6), Inches(5.8), code, 11)

# ============================================================
# SLIDE 8: auth_service.dart
# ============================================================
slide = title_slide('6. auth_service.dart - Services Firebase Auth')
tb = add_tb(slide, Inches(0.5), Inches(1.2), Inches(12), Inches(5.8), '', 14, BLACK)
tf = tb.text_frame; tf.word_wrap = True
tf.paragraphs[0].text = 'Rôle : Centralise toute la logique d\'authentification et de gestion des utilisateurs via Firebase'; tf.paragraphs[0].font.size = Pt(16); tf.paragraphs[0].font.color.rgb = GRAY
add_para(tf, '', 4)
add_para(tf, 'Méthodes Principales', 22, BLUE_DARK, True)

methods = [
    ('register(email, password, firstName, lastName, birthDate, phoneNumber)',
     'Crée un compte Firebase Auth + enregistre le profil dans Firestore (collection "users"). Calcul âge avec calculateAge(), bloque si < 13 ans.'),
    ('login(email, password)',
     'signInWithEmailAndPassword, puis getUserData() depuis Firestore. Si Firestore injoignable, crée un objet basique avec l\'email.'),
    ('logout()',
     'signOut() de Firebase Auth. Les SharedPreferences sont nettoyés par HomeScreen.'),
    ('resetPassword(email)',
     'sendPasswordResetEmail() de Firebase Auth. Gère les exceptions (utilisateur inexistant, etc.).'),
    ('getUserData(uid)',
     'Lit le document Firestore users/{uid}. Retourne Map avec firstName, lastName, email, etc.'),
    ('calculateAge(birthDate)',
     'Calcule l\'âge à partir de la date de naissance. Vérification côté client avant l\'inscription.'),
]
y = 2.0
for method, desc in methods:
    add_rounded_rect(slide, Inches(0.5), Inches(y), Inches(12.3), Inches(0.85), LIGHT_GRAY)
    add_tb(slide, Inches(0.8), Inches(y + 0.03), Inches(11.8), Inches(0.35), method, 13, BLUE_DARK, True)
    add_tb(slide, Inches(0.8), Inches(y + 0.38), Inches(11.8), Inches(0.45), desc, 12, GRAY)
    y += 0.95

add_para(tf, '', 6)
add_para(tf, 'Fallback Firestore', 18, BLUE_DARK, True)
add_bul(tf, 'Si Firestore est injoignable (réseau), login() retourne quand même des données basiques : {"firstName": email.split("@").first, "lastName": "", "email": email}', 13)
add_bul(tf, 'Évite de bloquer l\'utilisateur si le backend DB est temporairement hors ligne', 13)

# ============================================================
# SLIDE 9: login_screen.dart & register_screen.dart
# ============================================================
slide = title_slide('7. login_screen & register_screen - Authentification UI')
tb = add_tb(slide, Inches(0.5), Inches(1.2), Inches(6), Inches(5.8), '', 14, BLACK)
tf = tb.text_frame; tf.word_wrap = True
tf.paragraphs[0].text = 'login_screen.dart'; tf.paragraphs[0].font.size = Pt(22); tf.paragraphs[0].font.bold = True; tf.paragraphs[0].font.color.rgb = BLUE_DARK
add_bul(tf, 'Formulaire email + password avec validation (champs requis)', 14)
add_bul(tf, 'Bouton "Forgot Password?" -> Navigation vers ForgotScreen', 14)
add_bul(tf, 'Lien "Sign Up" -> Navigation vers RegisterScreen', 14)
add_bul(tf, 'Messages d\'erreur : container rouge #800000 avec bordure jaune (yellowAccent)', 14)
add_bul(tf, 'Gradient bleu, icône cadenas, Material 3 styling', 14)
add_bul(tf, 'Appel AuthService().login() avec try/catch', 14)
add_bul(tf, 'Callback onLogin(data) -> AppFlow -> HomeScreen', 14)

add_para(tf, '', 10)
add_para(tf, 'register_screen.dart', 22, BLUE_DARK, True)
add_bul(tf, '6 champs : First Name *, Last Name *, Email *, Password *, Confirm *, Phone', 14)
add_bul(tf, 'DatePicker pour Date of Birth *', 14)
add_bul(tf, 'Validation :', 14)
add_bul(tf, '  - Mot de passe >= 6 caractères', 13)
add_bul(tf, '  - Confirmation doit correspondre au mot de passe', 13)
add_bul(tf, '  - DatePicker obligatoire', 13)
add_bul(tf, '  - Âge >= 13 ans (AuthService.calculateAge)', 13)
add_bul(tf, 'Appel AuthService().register() avec tous les champs', 14)
add_bul(tf, 'SnackBar "Account created! Please login." + pop vers login', 14)

# Code snippet
code = '''// login_screen.dart - Appel Firebase
Future<void> _login() async {
  if (!_form.currentState!.validate()) return;
  setState(() { _loading = true; _error = ''; });
  try {
    final data = await AuthService().login(
      _emailCtrl.text.trim(), _passCtrl.text,
    );
    if (data != null && mounted) widget.onLogin(data);
    else setState(() => _error = 'Invalid credentials');
  } catch (e) {
    setState(() => _error = e.toString()
      .replaceAll('Exception: ', ''));
  } finally { if (mounted) setState(() => _loading = false); }
}

// register_screen.dart - Validation age
Future<void> _register() async {
  if (_auth.calculateAge(_birthDate!) < 13) {
    setState(() => _error = 'Must be 13 or older');
    return;
  }
  await _auth.register(
    email: _emailCtrl.text.trim(),
    password: _passCtrl.text,
    firstName: _fnameCtrl.text.trim(),
    lastName: _lnameCtrl.text.trim(),
    birthDate: _birthDate!,
    phoneNumber: _phoneCtrl.text.trim(),
  );
}'''
code_block(slide, Inches(6.8), Inches(1.2), Inches(6), Inches(5.8), code, 11)

# ============================================================
# SLIDE 10: home_screen.dart
# ============================================================
slide = title_slide('8. home_screen.dart - Navigation & Orchestration')
tb = add_tb(slide, Inches(0.5), Inches(1.2), Inches(6), Inches(5.8), '', 14, BLACK)
tf = tb.text_frame; tf.word_wrap = True
tf.paragraphs[0].text = 'Rôle : Écran principal après connexion, contient la navigation à 3 onglets'; tf.paragraphs[0].font.size = Pt(16); tf.paragraphs[0].font.color.rgb = GRAY
add_para(tf, '', 6)
add_para(tf, 'Structure', 22, BLUE_DARK, True)
add_bul(tf, 'StatefulWidget avec _idx (index de l\'onglet actif) et _audio (AudioService)', 14)
add_bul(tf, 'IndexedStack pour conserver l\'état de chaque onglet (pas de rebuild)', 14)
add_bul(tf, '3 enfants : StatsScreen, PlayerScreen, FavoritesScreen', 14)
add_bul(tf, 'NavigationBar Material 3 avec 3 destinations', 14)
add_para(tf, '', 6)
add_para(tf, 'AudioService Centralisé', 22, BLUE_DARK, True)
add_bul(tf, 'AudioService _audio = AudioService() créé DANS HomeScreen', 14)
add_bul(tf, 'Passé PAR CONSTRUCTEUR à PlayerScreen et FavoritesScreen', 14)
add_bul(tf, 'Évite les doublons de player et les fuites de listeners', 14)
add_bul(tf, '_audio.dispose() dans dispose() de HomeScreen', 14)
add_para(tf, '', 6)
add_para(tf, 'Logout', 22, BLUE_DARK, True)
add_bul(tf, 'Confirmation via AlertDialog "Are you sure?"', 14)
add_bul(tf, 'Nettoyage SharedPreferences : supprime clés listen_*, tracks_*, goal_hours', 14)
add_bul(tf, 'Appel AuthService().logout()', 14)
add_bul(tf, 'Navigation vers AppFlow() avec pushAndRemoveUntil (reset stack)', 14)

code = '''// home_screen.dart - Structure principale
class _HomeScreenState extends State<HomeScreen> {
  int _idx = 0;
  final AudioService _audio = AudioService(); // INSTANCE UNIQUE

  @override
  Widget build(BuildContext context) {
    return Scaffold(
      body: IndexedStack(index: _idx, children: [
        StatsScreen(userData: widget.userData),
        PlayerScreen(audio: _audio),     // injecté
        FavoritesScreen(audio: _audio,   // injecté
          onPlay: () => setState(() => _idx = 1)),
      ]),
      bottomNavigationBar: NavigationBar(
        selectedIndex: _idx,
        onDestinationSelected: (i) => setState(() => _idx = i),
        destinations: const [
          NavigationDestination(icon: Icon(Icons.bar_chart), label: 'Statistics'),
          NavigationDestination(icon: Icon(Icons.music_note), label: 'Player'),
          NavigationDestination(icon: Icon(Icons.favorite), label: 'Favorites'),
        ],
      ),
    );
  }
}'''
code_block(slide, Inches(6.8), Inches(1.2), Inches(6), Inches(5.8), code, 11)

# ============================================================
# SLIDE 11: player_screen.dart
# ============================================================
slide = title_slide('9. player_screen.dart - Lecteur Audio & Liste Sourates')
tb = add_tb(slide, Inches(0.5), Inches(1.2), Inches(6.2), Inches(5.8), '', 14, BLACK)
tf = tb.text_frame; tf.word_wrap = True
tf.paragraphs[0].text = 'Rôle : Affiche la liste des sourates et permet la lecture audio complète'; tf.paragraphs[0].font.size = Pt(16); tf.paragraphs[0].font.color.rgb = GRAY
add_para(tf, '', 6)
add_para(tf, 'Initialisation', 22, BLUE_DARK, True)
add_bul(tf, '_loadCategories() via QuranApiService.getCategories()', 14)
add_bul(tf, '_loadFavIds() via FavoriteService.stream().first', 14)
add_bul(tf, '4 listeners audio : onPositionChanged, onDurationChanged, onPlayerComplete, onPlayerStateChanged', 13)
add_para(tf, '', 6)
add_para(tf, 'Lecture d\'une Sourate', 22, BLUE_DARK, True)
add_bul(tf, '_playSurah(SurahCategory cat) :', 14)
add_bul(tf, '  1. Vérifie que cat.audioUrl n\'est pas null', 13)
add_bul(tf, '  2. Crée AudioTrack avec id, name, englishName, audioUrl', 13)
add_bul(tf, '  3. Reset _posSec et _durSec à 0, _playing = true', 13)
add_bul(tf, '  4. Appelle _audio.playTrack(track) avec try/catch', 13)
add_bul(tf, '  5. Si erreur -> _playing = false', 13)
add_para(tf, '', 6)
add_para(tf, 'Contrôles Audio', 22, BLUE_DARK, True)
add_bul(tf, 'Play/Pause : _audio.togglePlayPause()', 14)
add_bul(tf, 'Stop : _audio.stop() + reset position', 14)
add_bul(tf, 'Repeat : _audio.toggleRepeat() avec ReleaseMode.loop/stop', 14)
add_bul(tf, 'Slider : seek(Duration) avec clamp pour éviter les asserts', 14)
add_para(tf, '', 6)
add_para(tf, 'Gestion des Favoris', 22, BLUE_DARK, True)
add_bul(tf, '_toggleFav(AudioTrack) : Ajoute ou supprime des favoris', 14)
add_bul(tf, 'Suppression nécessite auth biométrique via BioService.requireAuthForDelete()', 14)
add_bul(tf, 'Sur web (kIsWeb) : pas de favoris ni biométrie', 14)
add_para(tf, '', 6)
add_para(tf, 'Liste des Sourates', 22, BLUE_DARK, True)
add_bul(tf, 'ListView.builder avec CircleAvatar (numéro), nom anglais/arabe', 14)
add_bul(tf, 'Sourate en cours : highlight bleu + icône play_arrow', 14)

code = '''Future<void> _playSurah(SurahCategory cat) async {
  final url = cat.audioUrl;
  if (url == null || url.isEmpty) return;
  final track = AudioTrack(id: cat.id,
    surahName: cat.name,
    surahEnglishName: cat.englishName,
    ayahNumber: 0, audioUrl: url);
  setState(() { _posSec = 0; _durSec = 0; _playing = true; });
  try {
    await _audio.playTrack(track);
  } catch (_) {
    if (mounted) setState(() => _playing = false);
  }
}'''
code_block(slide, Inches(6.8), Inches(1.2), Inches(6), Inches(3.5), code, 12)

# ============================================================
# SLIDE 12: stats_screen.dart
# ============================================================
slide = title_slide('10. stats_screen.dart - Statistiques d\'Écoute')
tb = add_tb(slide, Inches(0.5), Inches(1.2), Inches(6.2), Inches(5.8), '', 14, BLACK)
tf = tb.text_frame; tf.word_wrap = True
tf.paragraphs[0].text = 'Rôle : Affiche les statistiques d\'écoute de l\'utilisateur'; tf.paragraphs[0].font.size = Pt(16); tf.paragraphs[0].font.color.rgb = GRAY
add_para(tf, '', 6)
add_para(tf, 'Sections Affichées', 22, BLUE_DARK, True)
add_bul(tf, '1. Carte de Bienvenue : fond dégradé bleu, prénom + nom en gras (26pt)', 14)
add_bul(tf, '2. Temps Total : heures et minutes (ex: "12h 34m")', 14)
add_bul(tf, '3. Progression Objectif : LinearProgressIndicator (fl_chart), ratio total/objectif*60', 14)
add_bul(tf, '4. Histogramme 30 jours : fl_chart BarChart, ticks personnalisés, couleurs bleues', 14)
add_bul(tf, '5. Top Tracks : Liste des 5 sourates les plus écoutées', 14)
add_bul(tf, '6. Objectif Mensuel : Dropdown (5, 10, 15, 20, 25, 30, 40, 50h)', 14)
add_para(tf, '', 6)
add_para(tf, 'Sources de Données', 22, BLUE_DARK, True)
add_bul(tf, 'AudioService.getTotalMinutes() : Parcourt toutes les clés SharedPrefs listen_*', 14)
add_bul(tf, 'AudioService.getDailyMinutes() : Génère 30 entrées [date -> minutes]', 14)
add_bul(tf, 'AudioService.getTopTracks() : Agrège les clés tracks_*, trie par fréquence', 14)
add_bul(tf, 'AudioService.getMonthlyGoal() : Lit goal_hours des SharedPrefs (défaut 20)', 14)
add_para(tf, '', 6)
add_para(tf, 'Stockage Local', 22, BLUE_DARK, True)
add_bul(tf, 'listen_YYYY-MM-DD : Nombre de minutes écoutées par jour', 14)
add_bul(tf, 'tracks_YYYY-MM-DD : JSON array des IDs de sourates écoutées par jour', 14)
add_bul(tf, 'goal_hours : Objectif mensuel en heures', 14)
add_bul(tf, 'Nettoyé au logout dans HomeScreen._logout()', 14)

code = '''Future<void> _load() async {
  final total = await _audio.getTotalMinutes();
  final daily = await _audio.getDailyMinutes();
  final top = await _audio.getTopTracks();
  final goal = await _audio.getMonthlyGoal();
  setState(() { _totalMin = total; _daily = daily;
    _top = top; _goal = goal; _loading = false; });
}
// Histogramme fl_chart
Widget _buildChart() {
  return BarChart(BarChartData(
    barGroups: _daily.entries.map((e) =>
      BarChartGroupData(x: i++,
        barRods: [BarChartRodData(toY: e.value.toDouble(),
          color: Colors.blue)])).toList(),
  ));
}'''
code_block(slide, Inches(6.8), Inches(1.2), Inches(6), Inches(5.8), code, 12)

# ============================================================
# SLIDE 13: favorites_screen.dart
# ============================================================
slide = title_slide('11. favorites_screen.dart - Favoris Firestore')
tb = add_tb(slide, Inches(0.5), Inches(1.2), Inches(6), Inches(5.8), '', 14, BLACK)
tf = tb.text_frame; tf.word_wrap = True
tf.paragraphs[0].text = 'Rôle : Affiche la liste des favoris synchronisée en temps réel avec Firestore'; tf.paragraphs[0].font.size = Pt(16); tf.paragraphs[0].font.color.rgb = GRAY
add_para(tf, '', 6)
add_para(tf, 'Stream Firestore en Temps Réel', 22, BLUE_DARK, True)
add_bul(tf, 'StreamBuilder<QuerySnapshot> sur FavoriteService.stream()', 14)
add_bul(tf, 'stream() retourne FirebaseFirestore.collection("favorites").where("userId", ==, uid).snapshots()', 14)
add_bul(tf, 'Mise à jour automatique de l\'UI à chaque changement dans Firestore', 14)
add_para(tf, '', 6)
add_para(tf, 'Fonctionnalités', 22, BLUE_DARK, True)
add_bul(tf, 'Play : Appelle onPlay() callback -> HomeScreen bascule à l\'index Player', 14)
add_bul(tf, 'Delete : Vérifie isBiometricAvailable()', 14)
add_bul(tf, '  -> Si disponible : requireAuthForDelete() -> si OK -> FavoriteService.remove()', 13)
add_bul(tf, '  -> Sinon : suppression directe sans auth', 13)
add_bul(tf, 'Erreurs : catch + SnackBar avec message visible', 14)
add_para(tf, '', 6)
add_para(tf, 'Structure Firestore', 22, BLUE_DARK, True)
add_bul(tf, 'Collection : favorites', 14)
add_bul(tf, 'Document : {userId, surahId, surahName, surahEnglishName, audioUrl, addedAt}', 14)
add_bul(tf, 'Indexé par userId pour les requêtes', 14)

code = '''// Firestore stream
class FavoriteService {
  Stream<QuerySnapshot> stream() {
    final uid = FirebaseAuth.instance.currentUser!.uid;
    return FirebaseFirestore.instance
      .collection('favorites')
      .where('userId', isEqualTo: uid)
      .snapshots();
  }
}
// Delete with biometric check
Future<void> _delete(DocumentSnapshot doc) async {
  try {
    final available = await _bio!.isBiometricAvailable();
    if (available) {
      final ok = await _bio!.requireAuthForDelete();
      if (!ok) return;
    }
    await _fav!.remove(doc.id);
  } catch (e) {
    ScaffoldMessenger.of(context)
      .showSnackBar(SnackBar(content: Text('Error: $e')));
  }
}'''
code_block(slide, Inches(6.8), Inches(1.2), Inches(6), Inches(5.8), code, 11)

# ============================================================
# SLIDE 14: audio_service.dart
# ============================================================
slide = title_slide('12. audio_service.dart - Moteur Audio & Statistiques')
tb = add_tb(slide, Inches(0.5), Inches(1.2), Inches(6), Inches(5.8), '', 14, BLACK)
tf = tb.text_frame; tf.word_wrap = True
tf.paragraphs[0].text = 'Rôle : Centralise la lecture audio et l\'enregistrement des statistiques'; tf.paragraphs[0].font.size = Pt(16); tf.paragraphs[0].font.color.rgb = GRAY
add_para(tf, '', 6)
add_para(tf, 'AudioPlayer', 22, BLUE_DARK, True)
add_bul(tf, 'AudioPlayer() de audioplayers, stocké dans _player', 14)
add_bul(tf, 'playTrack(AudioTrack) : stop() -> play(UrlSource(url)) -> _recordListening()', 14)
add_bul(tf, 'togglePlayPause() : pause() ou resume() selon l\'état', 14)
add_bul(tf, 'toggleRepeat() : setReleaseMode(loop/stop)', 14)
add_bul(tf, 'seek(Duration) : positionnement dans le fichier audio', 14)
add_bul(tf, 'dispose() : libère le player (appelé par HomeScreen.dispose())', 14)
add_para(tf, '', 6)
add_para(tf, 'Enregistrement des Stats (_recordListening)', 22, BLUE_DARK, True)
add_bul(tf, 'Incrémente le compteur "listen_YYYY-MM-DD" dans SharedPreferences', 14)
add_bul(tf, 'Ajoute l\'ID de la piste à "tracks_YYYY-MM-DD" (JSON array)', 14)
add_bul(tf, 'Appelé automatiquement après chaque playTrack()', 14)
add_para(tf, '', 6)
add_para(tf, 'Lecture des Stats', 22, BLUE_DARK, True)
add_bul(tf, 'getTotalMinutes() : Somme de tous les listen_*', 14)
add_bul(tf, 'getDailyMinutes() : Map des 30 derniers jours [date -> minutes]', 14)
add_bul(tf, 'getTopTracks(limit: 5) : Agrège tracks_*, trie par occurrence', 14)
add_bul(tf, 'getMonthlyGoal() / setMonthlyGoal() : Lecture/écriture goal_hours', 14)
add_para(tf, '', 6)
add_para(tf, 'Gestion des Erreurs', 22, BLUE_DARK, True)
add_bul(tf, 'playTrack() catch + rethrow pour propagation au caller (PlayerScreen)', 14)
add_bul(tf, 'currentTrack = null en cas d\'erreur', 14)

code = '''class AudioService {
  final AudioPlayer _player = AudioPlayer();
  AudioTrack? _currentTrack;

  Future<void> playTrack(AudioTrack track) async {
    _currentTrack = track;
    try {
      await _player.stop();
      await _player.play(UrlSource(track.audioUrl));
      await _recordListening();
    } catch (e) {
      _currentTrack = null;
      rethrow;
    }
  }

  Future<void> _recordListening() async {
    final prefs = await SharedPreferences.getInstance();
    final today = DateTime.now().toString().split(' ')[0];
    final key = 'listen_$today';
    final current = prefs.getInt(key) ?? 0;
    await prefs.setInt(key, current + 1);
  }
}'''
code_block(slide, Inches(6.8), Inches(1.2), Inches(6), Inches(5.8), code, 11)

# ============================================================
# SLIDE 15: api_service.dart
# ============================================================
slide = title_slide('13. api_service.dart & Modèles')
tb = add_tb(slide, Inches(0.5), Inches(1.2), Inches(6), Inches(5.8), '', 14, BLACK)
tf = tb.text_frame; tf.word_wrap = True
tf.paragraphs[0].text = 'QuranApiService - Récupération des Sourates'; tf.paragraphs[0].font.size = Pt(22); tf.paragraphs[0].font.bold = True; tf.paragraphs[0].font.color.rgb = BLUE_DARK
add_bul(tf, 'Base URL : https://staticquran.vercel.app/api/v1', 14)
add_bul(tf, 'Endpoint : GET /surah -> retourne liste JSON des sourates', 14)
add_bul(tf, 'Parse la réponse : body["data"] -> pour chaque sourate :', 14)
add_bul(tf, '  - sequence -> id (string)', 13)
add_bul(tf, '  - name.arabic.short -> name (arabe)', 13)
add_bul(tf, '  - name.latin.short ou translation -> englishName', 13)
add_bul(tf, 'Génère audioUrl via _audioUrl(id) :', 14)
add_bul(tf, '  -> "https://server8.mp3quran.net/afs/{id3}.mp3" (id padded à 3 chiffres)', 13)
add_bul(tf, 'Fallback : 8 sourates locales si l\'API est injoignable', 14)

add_para(tf, '', 10)
add_para(tf, 'Modèles (track_model.dart)', 22, BLUE_DARK, True)
add_bul(tf, 'AudioTrack : id, surahName, surahEnglishName, ayahNumber, audioUrl', 14)
add_bul(tf, '  -> get title => \'$surahEnglishName - Ayah $ayahNumber\'', 13)
add_bul(tf, 'SurahCategory : id, name, englishName, audioUrl?, tracks[]', 14)
add_bul(tf, '  -> tracks est une liste d\'AudioTrack (vide pour une sourate complète)', 13)
add_bul(tf, 'Utilisés par : api_service (création), player_screen (affichage), audio_service (lecture)', 14)

code = '''class QuranApiService {
  static const String baseUrl =
    'https://staticquran.vercel.app/api/v1';

  String _audioUrl(String id) {
    final padded = id.padLeft(3, '0');
    return 'https://server8.mp3quran.net/afs/$padded.mp3';
  }

  Future<List<SurahCategory>> getCategories() async {
    try {
      final response = await http.get(
        Uri.parse('$baseUrl/surah'));
      if (response.statusCode == 200) {
        final body = jsonDecode(response.body);
        final List<dynamic> data = body['data'];
        return data.map((json) {
          final id = json['sequence'].toString();
          return SurahCategory(id: id,
            name: json['name']['arabic']['short'] ?? '',
            englishName: json['name']['latin']['short']
              ?? json['translation'] ?? '',
            audioUrl: _audioUrl(id), tracks: []);
        }).toList();
      }
    } catch (e) {}
    return _getFallbackCategories(); // 8 surahs
  }
}'''
code_block(slide, Inches(6.8), Inches(1.2), Inches(6), Inches(5.8), code, 11)

# ============================================================
# SLIDE 16: SECURITE
# ============================================================
slide = title_slide('14. Aspects Sécurité')
tb = add_tb(slide, Inches(0.5), Inches(1.2), Inches(12), Inches(5.8), '', 14, BLACK)
tf = tb.text_frame; tf.word_wrap = True
tf.paragraphs[0].text = 'Mécanismes de Sécurité Implémentés'; tf.paragraphs[0].font.size = Pt(24); tf.paragraphs[0].font.bold = True; tf.paragraphs[0].font.color.rgb = BLUE_DARK

sec_items = [
    ('Biométrie (local_auth)', 
     'AuthenticationOptions(biometricOnly: true) - Force l\'utilisation de l\'empreinte', 
     'isBiometricAvailable() vérifie que le capteur existe avant de proposer l\'auth',
     'Bouton Skip si indisponible -> ne bloque pas l\'utilisateur'),
    ('Firebase Auth',
     'Hachage et salage des mots de passe gérés par Firebase côté serveur',
     'Validation côté client : email valide, mot de passe >= 6 caractères',
     'Réinitialisation de mot de passe par email sécurisé Firebase'),
    ('Protection Favoris',
     'Suppression nécessite ré-authentification biométrique (requireAuthForDelete)',
     'Contournement si isBiometricAvailable() = false pour éviter le blocage',
     'SnackBar d\'erreur si la suppression échoue'),
    ('Permissions Android',
     'USE_BIOMETRIC + USE_FINGERPRINT : utilisation du capteur',
     'INTERNET + ACCESS_NETWORK_STATE : communication Firebase/API',
     'WAKE_LOCK + FOREGROUND_SERVICE : audio en arrière-plan'),
    ('Nettoyage Données',
     'Au logout : suppression des SharedPreferences (stats, objectifs)',
     'Évite les fuites de données entre utilisateurs sur le même appareil',
     'FirebaseAuth.signOut() invalide le token local'),
    ('Firestore Rules',
     'Règles de sécurité Firebase : accès limité à l\'UID de l\'utilisateur',
     'Collection "favorites" filtrée par userId côté serveur',
     'Pas d\'accès跨 utilisateur possible'),
]

y = 1.85
for title, d1, d2, d3 in sec_items:
    add_rounded_rect(slide, Inches(0.3), Inches(y), Inches(12.7), Inches(0.82), LIGHT_GRAY)
    add_tb(slide, Inches(0.5), Inches(y + 0.02), Inches(3), Inches(0.35), title, 16, BLUE_DARK, True)
    add_tb(slide, Inches(3.5), Inches(y + 0.02), Inches(9.5), Inches(0.25), d1, 12, GRAY)
    add_tb(slide, Inches(3.5), Inches(y + 0.28), Inches(9.5), Inches(0.25), d2, 12, GRAY)
    add_tb(slide, Inches(3.5), Inches(y + 0.54), Inches(9.5), Inches(0.25), d3, 12, GRAY)
    y += 0.88

# ============================================================
# SLIDE 17: DIFFICULTES
# ============================================================
slide = title_slide('Difficultés Rencontrées & Solutions')
tb = add_tb(slide, Inches(0.5), Inches(1.2), Inches(12), Inches(5.8), '', 14, BLACK)
tf = tb.text_frame; tf.word_wrap = True
tf.paragraphs[0].text = 'Problèmes Techniques et Leurs Résolutions'; tf.paragraphs[0].font.size = Pt(24); tf.paragraphs[0].font.bold = True; tf.paragraphs[0].font.color.rgb = BLUE_DARK

issues = [
    ('API Quran changée', 'yousefheiba.com retournait du HTML', 'Migration vers staticquran.vercel.app (JSON fiable)'),
    ('URL audio invalide', 'download.quranicaudio.com inaccessible', 'Migration vers server8.mp3quran.net/afs/{id}.mp3'),
    ('Slider assertion error', 'Plantage quand durée=0', 'Clamp de la valeur + maxVal=1.0 si dur=0'),
    ('AudioService dupliqué', 'Conflit de listeners audio', 'Instance unique dans HomeScreen, passée par constructeur'),
    ('Boucle de chargement', 'Spinner infini si user=null', 'setState explicite vers _showLogin=true'),
    ('Auto-login bypass', 'Mode démo ajouté par erreur', 'Suppression et retour au flux Firebase réel'),
    ('Emulateur Android 36', 'Activité introuvable (Error type 3)', 'Réinstallation + pm clear + attente boot'),
    ('Audio pas de son', 'pcmWrite I/O error driver ranchu', 'Test sur appareil physique nécessaire'),
    ('GPU instable', 'opengl32sw manquant', 'hw.gpu.mode=software dans config.ini'),
    ('Firestore injoignable', 'DNS/réseau sur cet ordinateur', 'Fallback dans login() avec données basiques'),
    ('Biométrie émulateur', 'Pas de capteur simulé par défaut', 'adb emu finger touch 1 + enregistrement préalable'),
]

y = 1.7
for prob, cause, sol in issues:
    add_rounded_rect(slide, Inches(0.3), Inches(y), Inches(12.7), Inches(0.6), LIGHT_GRAY)
    add_tb(slide, Inches(0.5), Inches(y + 0.02), Inches(2.5), Inches(0.35), prob, 14, BLUE_DARK, True)
    add_tb(slide, Inches(3.1), Inches(y + 0.02), Inches(4.2), Inches(0.55), cause, 12, GRAY)
    add_tb(slide, Inches(7.4), Inches(y + 0.02), Inches(5.5), Inches(0.55), sol, 12, GREEN)
    y += 0.65

# ============================================================
# SLIDE 18: CONCLUSION
# ============================================================
slide = new_slide()
add_bg(slide, BLUE_DARK)
add_shape(slide, 0, Inches(3.2), W, Inches(0.06), BLUE_LIGHT)
add_tb(slide, Inches(1), Inches(0.5), Inches(11), Inches(0.8), '15. Conclusion', 40, WHITE, True, PP_ALIGN.CENTER)

tb = add_tb(slide, Inches(1), Inches(1.5), Inches(11.3), Inches(5.5), '', 16, WHITE)
tf = tb.text_frame; tf.word_wrap = True
tf.paragraphs[0].text = 'Ce qui a été réalisé :'; tf.paragraphs[0].font.size = Pt(24); tf.paragraphs[0].font.bold = True; tf.paragraphs[0].font.color.rgb = WHITE
add_bul(tf, 'Application complète avec architecture MVC (Model-View-Controller)', 15)
add_bul(tf, 'Authentification biométrique robuste + Firebase Auth fonctionnelle', 15)
add_bul(tf, 'Lecture audio des 114 sourates avec API distante', 15)
add_bul(tf, 'Statistiques d\'écoute avec graphiques fl_chart et objectifs', 15)
add_bul(tf, 'Favoris synchronisés en temps réel via Firestore', 15)
add_bul(tf, 'Interface Material 3 responsive et personnalisée', 15)
add_bul(tf, 'Code sécurisé : permissions, biométrie, nettoyage données', 15)
add_para(tf, '', 10)
add_para(tf, 'Améliorations possibles :', 24, WHITE, True)
add_bul(tf, 'Lecture en arrière-plan (service Android foreground)', 15)
add_bul(tf, 'Téléchargement hors-ligne des sourates', 15)
add_bul(tf, 'Mode nuit et personnalisation du thème', 15)
add_bul(tf, 'Widget de rappel pour les prières', 15)
add_bul(tf, 'Version iOS complète', 15)
add_bul(tf, 'Tests unitaires et d\'intégration', 15)

# ============================================================
# SLIDE 19: MERCI
# ============================================================
slide = new_slide()
add_bg(slide, BLUE_DARK)
add_shape(slide, 0, Inches(3.2), W, Inches(0.06), BLUE_LIGHT)
add_tb(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.5), 'Merci de votre attention', 52, WHITE, True, PP_ALIGN.CENTER)
add_tb(slide, Inches(1), Inches(3.5), Inches(11), Inches(0.8), 'Des questions ?', 32, WHITE70, align=PP_ALIGN.CENTER)
add_rounded_rect(slide, Inches(3), Inches(5.0), Inches(7.3), Inches(0.8), WHITE)
add_tb(slide, Inches(3), Inches(5.1), Inches(7.3), Inches(0.6), 'github.com/dinapsdeslol/quranapp', 20, BLUE_DARK, True, PP_ALIGN.CENTER)

# Save
import os
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'Presentation_Quran_Player_V2.pptx')
prs.save(output_path)
print(f'Presentation saved to: {output_path}')
print(f'Total slides: {len(prs.slides)}')
